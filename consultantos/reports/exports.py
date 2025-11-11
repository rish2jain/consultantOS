"""
Export formats for reports (JSON, Excel, Word, PowerPoint)
"""
import json
import logging
from typing import Dict, Any, Optional
from io import BytesIO
from consultantos.models import StrategicReport

logger = logging.getLogger(__name__)


async def export_to_json(report: StrategicReport) -> Dict[str, Any]:
    """
    Export report as JSON

    Args:
        report: Strategic report to export

    Returns:
        Dictionary representation of the report
    """
    try:
        return report.model_dump()
    except Exception as e:
        logger.error(f"JSON export failed: {e}", exc_info=True)
        raise


async def export_to_powerpoint(report: StrategicReport) -> bytes:
    """
    Export report as PowerPoint presentation

    Args:
        report: Strategic report to export

    Returns:
        PowerPoint file as bytes
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Strategic Analysis Report"
        subtitle.text = (
            f"{report.executive_summary.company_name}\n"
            f"{report.executive_summary.industry}\n"
            f"Confidence Score: {report.executive_summary.confidence_score:.0%}"
        )

        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Executive Summary"

        tf = body_shape.text_frame
        tf.text = f"Company: {report.executive_summary.company_name}"

        p = tf.add_paragraph()
        p.text = f"Industry: {report.executive_summary.industry}"
        p.level = 0

        p = tf.add_paragraph()
        p.text = f"Confidence: {report.executive_summary.confidence_score:.0%}"
        p.level = 0

        # Slide 3: Key Findings
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Key Findings"

        tf = body_shape.text_frame
        tf.clear()

        for i, finding in enumerate(report.executive_summary.key_findings[:6]):
            if i == 0:
                tf.text = finding
            else:
                p = tf.add_paragraph()
                p.text = finding
                p.level = 0

        # Slide 4: Strategic Recommendation
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Strategic Recommendation"

        tf = body_shape.text_frame
        tf.text = report.executive_summary.strategic_recommendation

        # Save to bytes
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        raise ImportError("PowerPoint export requires python-pptx package")
    except Exception as e:
        logger.error(f"PowerPoint export failed: {e}", exc_info=True)
        raise


async def export_to_excel(report: StrategicReport) -> bytes:
    """
    Export report as Excel workbook
    
    Args:
        report: Strategic report to export
    
    Returns:
        Excel file as bytes
    """
    try:
        import pandas as pd
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Executive Summary"
        
        # Write executive summary
        ws.append(["Company", report.executive_summary.company_name])
        ws.append(["Industry", report.executive_summary.industry])
        ws.append(["Confidence Score", f"{report.executive_summary.confidence_score:.2%}"])
        ws.append([])
        ws.append(["Key Findings"])
        for finding in report.executive_summary.key_findings:
            ws.append([finding])
        
        ws.append([])
        ws.append(["Strategic Recommendation"])
        ws.append([report.executive_summary.strategic_recommendation])
        
        # Format headers
        for row in ws.iter_rows(min_row=1, max_row=10):
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.isupper():
                    cell.font = Font(bold=True)
        
        # Save to bytes
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        return output.getvalue()
    except ImportError:
        logger.error("openpyxl not installed. Install with: pip install openpyxl")
        raise ImportError("Excel export requires openpyxl package")
    except Exception as e:
        logger.error(f"Excel export failed: {e}", exc_info=True)
        raise


async def export_to_word(report: StrategicReport) -> bytes:
    """
    Export report as Word document
    
    Args:
        report: Strategic report to export
    
    Returns:
        Word document as bytes
    """
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Title
        title = doc.add_heading('Strategic Analysis Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Company info
        doc.add_heading('Company Information', level=1)
        doc.add_paragraph(f"Company: {report.executive_summary.company_name}")
        doc.add_paragraph(f"Industry: {report.executive_summary.industry}")
        doc.add_paragraph(f"Confidence Score: {report.executive_summary.confidence_score:.2%}")
        
        # Executive Summary
        doc.add_heading('Executive Summary', level=1)
        doc.add_heading('Key Findings', level=2)
        for finding in report.executive_summary.key_findings:
            doc.add_paragraph(finding, style='List Bullet')
        
        doc.add_heading('Strategic Recommendation', level=2)
        doc.add_paragraph(report.executive_summary.strategic_recommendation)
        
        # Save to bytes
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output.getvalue()
    except ImportError:
        logger.error("python-docx not installed. Install with: pip install python-docx")
        raise ImportError("Word export requires python-docx package")
    except Exception as e:
        logger.error(f"Word export failed: {e}", exc_info=True)
        raise