"""
Presentation generator for creating PowerPoint slides from narratives.
"""
import logging
from typing import List, Dict, Any, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None
    Inches = None
    Pt = None
    PP_ALIGN = None
    RGBColor = None

from consultantos.models.storytelling import (
    Narrative,
    Slide,
    Persona,
    PresentationRequest
)
from consultantos.storytelling.personas import get_persona_traits

logger = logging.getLogger(__name__)


class PresentationGenerator:
    """
    Generates PowerPoint presentations from narrative content.
    Adapts layout and design based on template and persona.
    """

    def __init__(self):
        """Initialize presentation generator."""
        if not HAS_PPTX:
            raise ImportError(
                "python-pptx package is required. "
                "Install with: pip install python-pptx"
            )
        self.default_width = Inches(10)
        self.default_height = Inches(7.5)

        # Brand color schemes
        self.color_schemes = {
            "professional": {
                "primary": RGBColor(41, 128, 185),    # Blue
                "secondary": RGBColor(52, 73, 94),    # Dark gray
                "accent": RGBColor(231, 76, 60),      # Red
                "background": RGBColor(255, 255, 255)  # White
            },
            "modern": {
                "primary": RGBColor(46, 204, 113),    # Green
                "secondary": RGBColor(52, 152, 219),  # Light blue
                "accent": RGBColor(155, 89, 182),     # Purple
                "background": RGBColor(236, 240, 241)  # Light gray
            },
            "minimal": {
                "primary": RGBColor(44, 62, 80),      # Dark blue
                "secondary": RGBColor(149, 165, 166), # Gray
                "accent": RGBColor(230, 126, 34),     # Orange
                "background": RGBColor(255, 255, 255)  # White
            },
            "corporate": {
                "primary": RGBColor(0, 59, 92),       # Navy
                "secondary": RGBColor(0, 123, 167),   # Blue
                "accent": RGBColor(227, 114, 34),     # Orange
                "background": RGBColor(255, 255, 255)  # White
            }
        }

    async def generate_presentation(
        self,
        narrative: Narrative,
        max_slides: int = 10,
        template: str = "professional",
        brand_colors: Optional[Dict[str, str]] = None,
        include_charts: bool = True
    ) -> tuple[Presentation, List[Slide]]:
        """
        Generate PowerPoint presentation from narrative.

        Args:
            narrative: Narrative content to convert
            max_slides: Maximum number of slides
            template: Template style (professional, modern, minimal, corporate)
            brand_colors: Optional custom brand colors
            include_charts: Whether to include chart placeholders

        Returns:
            Tuple of (Presentation object, Slide metadata list)
        """
        logger.info(f"Generating presentation with {max_slides} max slides using {template} template")

        # Create presentation
        prs = Presentation()
        prs.slide_width = self.default_width
        prs.slide_height = self.default_height

        # Get color scheme
        colors = self._get_color_scheme(template, brand_colors)

        # Generate slide metadata
        slide_metadata = self._generate_slide_structure(
            narrative=narrative,
            max_slides=max_slides,
            include_charts=include_charts
        )

        # Create slides
        for slide_meta in slide_metadata:
            self._add_slide(
                prs=prs,
                slide_meta=slide_meta,
                colors=colors,
                persona=narrative.generated_for_persona
            )

        logger.info(f"Generated presentation with {len(prs.slides)} slides")
        return prs, slide_metadata

    def _generate_slide_structure(
        self,
        narrative: Narrative,
        max_slides: int,
        include_charts: bool
    ) -> List[Slide]:
        """Generate slide structure from narrative."""
        slides = []
        slide_num = 1

        # Title slide
        slides.append(Slide(
            slide_number=slide_num,
            layout="title",
            title=narrative.title,
            content=[narrative.subtitle] if narrative.subtitle else [],
            speaker_notes=f"Introduction to {narrative.title}"
        ))
        slide_num += 1

        # Key insights slide
        if narrative.key_insights and slide_num <= max_slides:
            slides.append(Slide(
                slide_number=slide_num,
                layout="title_content",
                title="Key Insights",
                content=narrative.key_insights[:5],
                speaker_notes="Overview of main findings"
            ))
            slide_num += 1

        # Section slides
        for section in narrative.sections:
            if slide_num > max_slides:
                break

            # Section content slide
            content_items = section.key_points if section.key_points else [section.content[:200] + "..."]

            slide = Slide(
                slide_number=slide_num,
                layout="title_content" if not section.visualizations else "two_column",
                title=section.heading,
                content=content_items[:6],  # Max 6 bullets
                chart_ids=section.visualizations if include_charts else [],
                speaker_notes=section.content[:500]
            )
            slides.append(slide)
            slide_num += 1

        # Recommendations slide
        if narrative.recommendations and slide_num <= max_slides:
            slides.append(Slide(
                slide_number=slide_num,
                layout="title_content",
                title="Recommendations",
                content=narrative.recommendations[:5],
                speaker_notes="Actionable next steps"
            ))
            slide_num += 1

        # Closing slide
        if slide_num <= max_slides:
            slides.append(Slide(
                slide_number=slide_num,
                layout="title",
                title="Thank You",
                content=["Questions?"],
                speaker_notes="Summary and Q&A"
            ))

        return slides[:max_slides]

    def _add_slide(
        self,
        prs: Presentation,
        slide_meta: Slide,
        colors: Dict[str, RGBColor],
        persona: Persona
    ) -> None:
        """Add a slide to the presentation."""
        # Choose layout
        if slide_meta.layout == "title":
            slide_layout = prs.slide_layouts[0]  # Title slide
        elif slide_meta.layout == "title_content":
            slide_layout = prs.slide_layouts[1]  # Title and content
        elif slide_meta.layout == "two_column":
            slide_layout = prs.slide_layouts[3]  # Two content
        else:
            slide_layout = prs.slide_layouts[5]  # Blank

        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background"]

        # Add title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_meta.title
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]

        # Add content based on layout
        if slide_meta.layout == "title_content":
            self._add_content_bullets(slide, slide_meta.content, colors)

        elif slide_meta.layout == "two_column":
            self._add_two_column_content(slide, slide_meta, colors)

        # Add speaker notes
        if slide_meta.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_meta.speaker_notes

    def _add_content_bullets(
        self,
        slide,
        content: List[str],
        colors: Dict[str, RGBColor]
    ) -> None:
        """Add bullet point content to slide."""
        # Find content placeholder
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                text_frame.clear()

                for item in content:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors["secondary"]
                    p.space_before = Pt(12)

                break

    def _add_two_column_content(
        self,
        slide,
        slide_meta: Slide,
        colors: Dict[str, RGBColor]
    ) -> None:
        """Add two-column content (text + chart placeholder)."""
        # Left column: bullet points
        left = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.5), Inches(5)
        )

        text_frame = left.text_frame
        for item in slide_meta.content[:4]:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = colors["secondary"]

        # Right column: chart placeholder
        if slide_meta.chart_ids:
            chart_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(1.5),
                Inches(4), Inches(5)
            )
            chart_frame = chart_box.text_frame
            p = chart_frame.paragraphs[0]
            p.text = f"[Chart: {slide_meta.chart_ids[0]}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = colors["accent"]

    def _get_color_scheme(
        self,
        template: str,
        brand_colors: Optional[Dict[str, str]]
    ) -> Dict[str, RGBColor]:
        """Get color scheme for presentation."""
        if brand_colors:
            # Convert hex colors to RGB
            return {
                key: self._hex_to_rgb(value)
                for key, value in brand_colors.items()
            }

        return self.color_schemes.get(template, self.color_schemes["professional"])

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    async def export_presentation(
        self,
        prs: Presentation,
        output_path: str
    ) -> str:
        """
        Export presentation to file.

        Args:
            prs: Presentation object
            output_path: Output file path

        Returns:
            Path to saved file
        """
        try:
            prs.save(output_path)
            logger.info(f"Saved presentation to {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Failed to save presentation: {str(e)}")
            raise

    async def generate_slide_preview(
        self,
        slide_meta: Slide
    ) -> Dict[str, Any]:
        """
        Generate a preview representation of a slide.

        Args:
            slide_meta: Slide metadata

        Returns:
            Preview data dictionary
        """
        return {
            "slide_number": slide_meta.slide_number,
            "layout": slide_meta.layout,
            "title": slide_meta.title,
            "content_items": len(slide_meta.content),
            "has_charts": len(slide_meta.chart_ids) > 0,
            "chart_count": len(slide_meta.chart_ids)
        }
