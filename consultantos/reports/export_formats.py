"""
Export Formats for Enhanced Reports

Supports exporting enhanced reports to:
- JSON (structured data)
- Excel (spreadsheets with multiple sheets)
- Word (editable documents)
"""
import json
from typing import Dict, Any, Optional
from datetime import datetime
from consultantos.models.enhanced_reports import EnhancedStrategicReport
import logging

logger = logging.getLogger(__name__)


def export_to_json(enhanced_report: EnhancedStrategicReport) -> str:
    """
    Export enhanced report to JSON format.
    
    Args:
        enhanced_report: Enhanced strategic report
        
    Returns:
        JSON string
    """
    try:
        # Convert to dict and serialize
        report_dict = enhanced_report.model_dump()
        return json.dumps(report_dict, indent=2, default=str)
    except Exception as e:
        logger.error(f"Failed to export to JSON: {e}", exc_info=True)
        raise


def export_to_excel(enhanced_report: EnhancedStrategicReport) -> bytes:
    """
    Export enhanced report to Excel format.
    
    Creates multiple sheets:
    - Executive Summary
    - Recommendations
    - Risks
    - Opportunities
    - Framework Analysis
    
    Args:
        enhanced_report: Enhanced strategic report
        
    Returns:
        Excel file bytes
    """
    try:
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
        except ImportError:
            raise ImportError(
                "openpyxl is required for Excel export. "
                "Install it with: pip install openpyxl. "
                "Alternatively, use export_to_json() for JSON format."
            )
        
        # Create workbook
        wb = openpyxl.Workbook()
        
        # Remove default sheet
        if 'Sheet' in wb.sheetnames:
            wb.remove(wb['Sheet'])
        
        # Executive Summary Sheet
        ws_summary = wb.create_sheet("Executive Summary")
        ws_summary.append(["ENHANCED STRATEGIC ANALYSIS REPORT"])
        ws_summary.append([])
        
        exec_summary = enhanced_report.executive_summary_layer
        ws_summary.append(["Company", exec_summary.company_overview.get("name", "N/A")])
        ws_summary.append(["Industry", exec_summary.company_overview.get("industry", "N/A")])
        ws_summary.append(["Analysis Date", exec_summary.analysis_date.strftime("%Y-%m-%d")])
        ws_summary.append(["Confidence Score", f"{exec_summary.confidence_score:.1%}"])
        ws_summary.append([])
        
        ws_summary.append(["Key Findings"])
        for finding in exec_summary.key_findings:
            ws_summary.append([f"• {finding}"])
        ws_summary.append([])
        
        ws_summary.append(["Strategic Recommendations"])
        for rec in exec_summary.strategic_recommendations:
            ws_summary.append([f"• {rec}"])
        
        # Recommendations Sheet
        ws_recs = wb.create_sheet("Recommendations")
        ws_recs.append(["Title", "Priority", "Timeline", "Owner", "Expected Outcome", "Success Metrics"])
        
        recommendations = enhanced_report.actionable_recommendations
        all_actions = (
            recommendations.immediate_actions +
            recommendations.short_term_actions +
            recommendations.medium_term_actions +
            recommendations.long_term_actions
        )
        
        for action in all_actions:
            ws_recs.append([
                action.title,
                action.priority.value,
                action.timeline.value,
                action.owner or "Unassigned",
                action.expected_outcome,
                "; ".join(action.success_metrics)
            ])
        
        # Risks Sheet
        ws_risks = wb.create_sheet("Risks")
        ws_risks.append(["Title", "Description", "Likelihood", "Impact", "Risk Score", "Mitigation Strategies"])
        
        for risk in enhanced_report.risk_opportunity_matrix.risks:
            ws_risks.append([
                risk.title,
                risk.description,
                risk.likelihood,
                risk.impact.value,
                risk.risk_score,
                "; ".join(risk.mitigation_strategies)
            ])
        
        # Opportunities Sheet
        ws_opps = wb.create_sheet("Opportunities")
        ws_opps.append(["Title", "Description", "Impact Potential", "Feasibility", "Priority Score", "Timeline"])
        
        for opp in enhanced_report.risk_opportunity_matrix.opportunities:
            ws_opps.append([
                opp.title,
                opp.description,
                opp.impact_potential,
                opp.feasibility,
                opp.priority_score,
                f"{opp.timeline_to_value} months"
            ])
        
        # Save to bytes
        from io import BytesIO
        excel_buffer = BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        return excel_buffer.read()
    
    except Exception as e:
        logger.error(f"Failed to export to Excel: {e}", exc_info=True)
        # Fallback to JSON
        return export_to_json(enhanced_report).encode('utf-8')


def export_to_powerpoint(enhanced_report: EnhancedStrategicReport) -> bytes:
    """
    Export enhanced report to PowerPoint format.

    Creates executive presentation with:
    - Executive Summary slide
    - Key Findings slides
    - Strategic Recommendations
    - Risk/Opportunity Matrix
    - Action Plan

    Args:
        enhanced_report: Enhanced strategic report

    Returns:
        PowerPoint file bytes
    """
    try:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint export. "
                "Install it with: pip install python-pptx. "
                "Alternatively, use export_to_json() for JSON format."
            )

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        exec_summary = enhanced_report.executive_summary_layer

        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Strategic Analysis Report"
        subtitle.text = (
            f"{exec_summary.company_overview.get('name', 'Company')}\n"
            f"{exec_summary.analysis_date.strftime('%B %d, %Y')}\n"
            f"Confidence Score: {exec_summary.confidence_score:.0%}"
        )

        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Executive Summary"

        tf = body_shape.text_frame
        tf.text = f"Company: {exec_summary.company_overview.get('name', 'N/A')}"

        p = tf.add_paragraph()
        p.text = f"Industry: {exec_summary.company_overview.get('industry', 'N/A')}"
        p.level = 0

        p = tf.add_paragraph()
        p.text = f"Analysis Confidence: {exec_summary.confidence_score:.0%}"
        p.level = 0

        # Slide 3: Key Findings
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Key Findings"

        tf = body_shape.text_frame
        tf.clear()

        for i, finding in enumerate(exec_summary.key_findings[:5]):  # Top 5
            if i == 0:
                tf.text = finding
            else:
                p = tf.add_paragraph()
                p.text = finding
                p.level = 0

        # Slide 4: Strategic Recommendations
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Strategic Recommendations"

        tf = body_shape.text_frame
        tf.clear()

        for i, rec in enumerate(exec_summary.strategic_recommendations[:5]):  # Top 5
            if i == 0:
                tf.text = rec
            else:
                p = tf.add_paragraph()
                p.text = rec
                p.level = 0

        # Slide 5: Top Risks
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Top Risks"

        tf = body_shape.text_frame
        tf.clear()

        risk_opp = enhanced_report.risk_opportunity_matrix
        for i, risk in enumerate(risk_opp.risks[:4]):  # Top 4
            risk_text = f"{risk.title} (Risk Score: {risk.risk_score}/10)"
            if i == 0:
                tf.text = risk_text
            else:
                p = tf.add_paragraph()
                p.text = risk_text
                p.level = 0

            # Add mitigation strategy
            if risk.mitigation_strategies:
                p = tf.add_paragraph()
                p.text = f"Mitigation: {risk.mitigation_strategies[0]}"
                p.level = 1

        # Slide 6: Top Opportunities
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Top Opportunities"

        tf = body_shape.text_frame
        tf.clear()

        for i, opp in enumerate(risk_opp.opportunities[:4]):  # Top 4
            opp_text = f"{opp.title} (Impact: {opp.impact_potential}/10, Priority: {opp.priority_score}/10)"
            if i == 0:
                tf.text = opp_text
            else:
                p = tf.add_paragraph()
                p.text = opp_text
                p.level = 0

        # Slide 7: Action Plan
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Immediate Action Plan"

        tf = body_shape.text_frame
        tf.clear()

        recommendations = enhanced_report.actionable_recommendations
        immediate_actions = (
            recommendations.critical_actions[:2] +
            recommendations.immediate_actions[:3]
        )

        for i, action in enumerate(immediate_actions[:5]):  # Top 5 actions
            action_text = f"{action.title} ({action.timeline.value})"
            if i == 0:
                tf.text = action_text
            else:
                p = tf.add_paragraph()
                p.text = action_text
                p.level = 0

            # Add owner
            p = tf.add_paragraph()
            p.text = f"Owner: {action.owner or 'Unassigned'}"
            p.level = 1

        # Save to bytes
        from io import BytesIO
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.read()

    except Exception as e:
        logger.error(f"Failed to export to PowerPoint: {e}", exc_info=True)
        # Fallback to JSON
        return export_to_json(enhanced_report).encode('utf-8')


def export_to_word(enhanced_report: EnhancedStrategicReport) -> bytes:
    """
    Export enhanced report to Word format.

    Args:
        enhanced_report: Enhanced strategic report

    Returns:
        Word document bytes
    """
    try:
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX export. "
                "Install it with: pip install python-docx. "
                "Alternatively, use export_to_json() for JSON format."
            )
        
        # Create document
        doc = Document()
        
        # Title
        title = doc.add_heading('Enhanced Strategic Analysis Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Executive Summary
        doc.add_heading('Executive Summary', 1)
        
        exec_summary = enhanced_report.executive_summary_layer
        doc.add_paragraph(f"Company: {exec_summary.company_overview.get('name', 'N/A')}")
        doc.add_paragraph(f"Industry: {exec_summary.company_overview.get('industry', 'N/A')}")
        doc.add_paragraph(f"Analysis Date: {exec_summary.analysis_date.strftime('%Y-%m-%d')}")
        doc.add_paragraph(f"Confidence Score: {exec_summary.confidence_score:.1%}")
        
        doc.add_heading('Key Findings', 2)
        for finding in exec_summary.key_findings:
            doc.add_paragraph(finding, style='List Bullet')
        
        doc.add_heading('Strategic Recommendations', 2)
        for rec in exec_summary.strategic_recommendations:
            doc.add_paragraph(rec, style='List Bullet')
        
        # Actionable Recommendations
        doc.add_heading('Actionable Recommendations', 1)
        
        recommendations = enhanced_report.actionable_recommendations
        if recommendations.critical_actions:
            doc.add_heading('Critical Actions', 2)
            for action in recommendations.critical_actions:
                p = doc.add_paragraph()
                p.add_run(f"{action.title}").bold = True
                doc.add_paragraph(f"Priority: {action.priority.value} | Timeline: {action.timeline.value}")
                doc.add_paragraph(f"Owner: {action.owner or 'Unassigned'}")
                doc.add_paragraph(f"Expected Outcome: {action.expected_outcome}")
        
        # Risk & Opportunity
        doc.add_heading('Risk & Opportunity Assessment', 1)
        
        risk_opp = enhanced_report.risk_opportunity_matrix
        
        doc.add_heading('Top Risks', 2)
        for risk in risk_opp.risks[:5]:
            p = doc.add_paragraph()
            p.add_run(f"{risk.title}").bold = True
            doc.add_paragraph(f"Likelihood: {risk.likelihood}/10 | Impact: {risk.impact.value} | Risk Score: {risk.risk_score}/10")
        
        doc.add_heading('Top Opportunities', 2)
        for opp in risk_opp.opportunities[:5]:
            p = doc.add_paragraph()
            p.add_run(f"{opp.title}").bold = True
            doc.add_paragraph(f"Impact: {opp.impact_potential}/10 | Feasibility: {opp.feasibility}/10 | Priority: {opp.priority_score}/10")
        
        # Save to bytes
        from io import BytesIO
        word_buffer = BytesIO()
        doc.save(word_buffer)
        word_buffer.seek(0)
        return word_buffer.read()
    
    except Exception as e:
        logger.error(f"Failed to export to Word: {e}", exc_info=True)
        # Fallback to JSON
        return export_to_json(enhanced_report).encode('utf-8')

